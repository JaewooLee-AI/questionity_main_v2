import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_manual_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # Ace Hotel Editorial Color Tokens
    BG_WARM = RGBColor(244, 243, 238)     # #F4F3EE Paper Off-White
    BG_DARK = RGBColor(26, 26, 26)        # #1A1A1A Matte Black
    BG_CHARCOAL = RGBColor(36, 36, 34)    # #242422 Dark Charcoal Gray
    BG_WHITE = RGBColor(255, 255, 255)
    TEXT_BLACK = RGBColor(26, 26, 26)     # Deep Black
    TEXT_WHITE = RGBColor(244, 243, 238)   # Off-White
    TEXT_MUTED = RGBColor(100, 100, 100)
    COLOR_BRICK = RGBColor(140, 35, 24)   # #8C2318 Ace Brick Red
    COLOR_ORANGE = RGBColor(255, 100, 51)  # #FF6433 Vivid Orange Accent
    BORDER_GRAY = RGBColor(200, 195, 185)

    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_run(paragraph, text, size_pt=12, bold=False, color=TEXT_BLACK, font_name="Arial"):
        run = paragraph.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color
        return run

    def add_header(slide, title_text, category_text="QUESTIONITY MANUAL"):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        add_run(p0, category_text.upper(), size_pt=11, bold=True, color=COLOR_BRICK)

        p1 = tf.add_paragraph()
        add_run(p1, title_text, size_pt=22, bold=True, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 1: COVER
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide1, BG_DARK)

    # Outer Brick Red Border Box
    frame = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6), Inches(12.133), Inches(6.3))
    frame.fill.background()
    frame.line.color.rgb = COLOR_BRICK
    frame.line.width = Pt(2.5)

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(10.933), Inches(4.7))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    add_run(p, "QUESTION + COMMUNITY PLATFORM\n", size_pt=13, bold=True, color=COLOR_BRICK)

    p = tf1.add_paragraph()
    add_run(p, "퀘스처니티 (Questionity)\n", size_pt=38, bold=True, color=TEXT_WHITE)

    p = tf1.add_paragraph()
    add_run(p, "사용자 및 운영자 통합 매뉴얼 (User & Admin Guidebook)\n\n", size_pt=22, bold=False, color=RGBColor(210, 210, 210))

    p = tf1.add_paragraph()
    add_run(p, "• 플랫폼 특장점: ", size_pt=13, bold=True, color=COLOR_BRICK)
    add_run(p, "에이스 호텔 에디토리얼 디자인, 300권 도서 큐레이션, Streamlit AI 자동화 & Supabase 인프라\n", size_pt=13, bold=False, color=RGBColor(190, 190, 190))

    p = tf1.add_paragraph()
    add_run(p, "• 대상: ", size_pt=13, bold=True, color=COLOR_BRICK)
    add_run(p, "일반 모임 참여 사용자 & 독서클럽 운영 관리자 (Admin)\n", size_pt=13, bold=False, color=RGBColor(190, 190, 190))

    p = tf1.add_paragraph()
    add_run(p, "• 최종 수정일: ", size_pt=12, bold=True, color=COLOR_BRICK)
    add_run(p, "2026년 8월 20일   |   ", size_pt=12, bold=False, color=TEXT_WHITE)
    add_run(p, "• 배포 버전: ", size_pt=12, bold=True, color=COLOR_BRICK)
    add_run(p, "v2.0 Production Ready", size_pt=12, bold=False, color=TEXT_WHITE)

    # =========================================================================
    # SLIDE 2: TABLE OF CONTENTS
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide2, BG_WARM)
    add_header(slide2, "목차 (Table of Contents)")

    toc_data = [
        ("01", "브랜드 비전 & 에이스 호텔 디자인 시스템", "혜화 창경궁로 독서 아지트 포지셔닝 및 브루탈리즘 에디토리얼 UI 컨셉"),
        ("02", "사용자 매뉴얼 - 상단 헤더 & 모집중 독서모임", "에이스 상단 SIGN UP 버튼 셀, 모집중 차콜 그레이 100% Full-Bleed 배너"),
        ("03", "사용자 매뉴얼 - 생생 후기 & 다중 후기 모달", "Full-Bleed 수동 화살표 포스터 카드, READ MORE 클릭 시 책별 전체 후기 조망"),
        ("04", "사용자 매뉴얼 - 쿠폰 패스 & 오프라인 라운지", "2·6·10장 차감 패스 가이드, 혜화 창경궁로 270 길찾기 및 FAQ/1:1 문의"),
        ("05", "운영자 매뉴얼 - Streamlit AI 대시보드", "알라딘 TTB 키 연동, Gemini AI 350자+ 리뷰 생성, Supabase 원클릭 출판"),
        ("06", "시스템 인프라 & Vercel ENV 설정 가이드", "Vercel 빌드 설정(out), Supabase DB 테이블 구조 및 Vercel 환경변수 (.env)")
    ]

    for idx, (num, title, desc) in enumerate(toc_data):
        col = idx % 2
        row = idx // 3
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.7 + row * 1.7)

        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.6), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.25)
        tf_c.margin_top = Inches(0.2)

        p = tf_c.paragraphs[0]
        add_run(p, f"SECTION {num}\n", size_pt=10, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, title, size_pt=15, bold=True, color=TEXT_BLACK)

        p = tf_c.add_paragraph()
        add_run(p, desc, size_pt=11, bold=False, color=TEXT_MUTED)

    # =========================================================================
    # SLIDE 3: BRAND VISION & DESIGN SYSTEM
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide3, BG_WARM)
    add_header(slide3, "01. 브랜드 비전 & 에이스 호텔 디자인 시스템")

    concepts = [
        ("에이스 호텔 디자인 시스템", [
            "• 100% 직사각형 레이아웃 (.rounded-none)",
            "• 웜 캔버스 (#F4F3EE) & 딥 차콜 (#242422)",
            "• 에이스 시그니처 브릭 레드 (#8C2318) 포인트",
            "• Gmarket Sans / Bebas Neue 타이포그래피"
        ]),
        ("스타트업 포지셔닝", [
            "• 혜화 창경궁로 270 오프라인 독서 아지트",
            "• 진정성 있는 초기 스타트업 커뮤니티",
            "• 1st Season 창립 멤버십 운영",
            "• 질문 기반의 소그룹 액티브 토론"
        ]),
        ("300권 검증 도서 큐레이션", [
            "• 알라딘 베스트셀러 50권씩 6개 트랙",
            "• 인문/사회/철학, 비즈니스/재테크,",
            "  과학/기술, 문학/예술, 자기계발/심리,",
            "  스타트업/디자인 명작 클래식 포함"
        ])
    ]

    for idx, (title, items) in enumerate(concepts):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n주요 가이드라인:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 4: USER MANUAL 1 - HEADER & RECRUITING CLUBS
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide4, BG_WARM)
    add_header(slide4, "02. 사용자 매뉴얼 - 상단 헤더 & 모집중인 독서모임")

    user_h_clubs = [
        ("상단 헤더 SIGN UP 셀 UI", [
            "• 위치: 상단 네비게이션 우측 끝 독점 셀",
            "• UI: 에이스 호텔 BOOK NOW 100% 동일 색상/UI",
            "• 브릭 레드(#8C2318) & 호버 시 매트 블랙(#1A1A1A)",
            "• 비로그인 시 SIGN UP (/signup), 로그인 시 MY PAGE"
        ]),
        ("모집중 독서모임 차콜 배너", [
            "• 모집중인 독서모임 (5 CLUBS RECRUITING)",
            "• 화면 좌우 100% Full-Bleed 차콜 그레이 (#242422)",
            "• 모집중 전용 딥 차콜 아우라 카드 렌더링",
            "• 솔리드 하이콘트라스트 [상세내용] 버튼 제공"
        ]),
        ("진행중/종료 트랙 레이아웃", [
            "• ⚡ 진행중인 독서모임 (CLUBS ACTIVE)",
            "• ✅ 종료된 독서모임 (CLUBS ARCHIVED)",
            "• 페이퍼 오프화이트 정갈한 카드로 교차 배치",
            "• 원형 화살표 휠 이동 지원"
        ])
    ]

    for idx, (title, items) in enumerate(user_h_clubs):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n기능 사양:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 5: USER MANUAL 2 - REVIEWS & MULTI-REVIEW MODAL
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide5, BG_WARM)
    add_header(slide5, "03. 사용자 매뉴얼 - 생생 후기 & 다중 후기 스택 모달")

    reviews_spec = [
        ("Full-Bleed 포스터 카드", [
            "• 검은 테두리와 여백 0% 제거 (Full-Bleed)",
            "• aspect-[4/5] 책 표지 100% 꽉 차게 렌더링",
            "• 외곽선 없는 깔끔한 페이퍼 카드 룩",
            "• 20개 AI 생생 후기 전량 100% 로딩"
        ]),
        ("수동 원형 화살표 슬라이더", [
            "• 마키(Marquee) 자동 흐름 완전 제거",
            "• 좌우 플로팅 원형 화살표 버튼 (← / →)",
            "• 사용자가 원하는 속도로 직접 조작",
            "• 터치/드래그 스크롤 완벽 지원"
        ]),
        ("다중 후기 스택 모달 팝업", [
            "• READ MORE (2) 클릭 시 도서별 통합 모달",
            "• 동일 도서에 대한 모든 멤버 후기 모아보기",
            "• 작성자 페르소나, 별점, 후기 전문 조망",
            "• 개별 LIKE (좋아요) 카운터 탑재"
        ])
    ]

    for idx, (title, items) in enumerate(reviews_spec):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n기능 사양:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 6: USER MANUAL 3 - PASS GUIDE & SEOUL LOUNGE
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide6, BG_WARM)
    add_header(slide6, "04. 사용자 매뉴얼 - 쿠폰 패스 & 오프라인 라운지")

    user_pass_loc = [
        ("멤버십 패스 가이드 (#how-it-works)", [
            "• 3단계 서비스 프로토콜 (탐색 -> 패스 충전 -> 차감)",
            "• 2장 패스 (기본), 6장 패스 (15% 할인 추천), 10장 패스 (20%)",
            "• 독서클럽 신청 시 간편하게 [1장씩 차감] 신청",
            "• 반응형 줄바꿈 (break-keep) 완벽 적용"
        ]),
        ("오프라인 라운지 티켓 (#location)", [
            "• 혜화 독서 아지트: 서울특별시 종로구 창경궁로 270",
            "• 오프라인 라운지 룸 101 사진 배너",
            "• [도로명 주소 복사] & [네이버 지도 길찾기] 원클릭",
            "• 4가지 아메니티 (Free Coffee, Vinyl Lounge 등)"
        ])
    ]

    for idx, (title, items) in enumerate(user_pass_loc):
        left = Inches(0.8 + idx * 5.9)
        top = Inches(1.7)
        card = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.6), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.3)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=15, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n사용자 이용 안내:\n\n", size_pt=12, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=12, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 7: USER MANUAL 4 - FAQ & 1:1 INQUIRY SYSTEM
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide7, BG_WARM)
    add_header(slide7, "04. 사용자 매뉴얼 - FAQ & 1:1 문의 접수 (/faq)")

    faq_inquiry = [
        ("FAQ 아코디언 & 실시간 검색", [
            "• 검색바: 질문 및 키워드 입력 시 실시간 필터링",
            "• 카테고리 탭: [모임신청], [결제/패스], [장소/시설], [환불규정]",
            "• 아코디언: 클릭 시 깔끔한 에이스 호텔 스타일 답변 토글"
        ]),
        ("1:1 문의하기 폼 (Form Submission)", [
            "• 수집 데이터: 성함, 이메일, 연락처, 문의 유형, 제목, 내용",
            "• 유효성 검사: 필수 입력란 미작성 시 실시간 피드백",
            "• 접수 결과: Supabase `inquiries` 데이터베이스에 즉시 저장",
            "• 상태값: `pending` (접수 완료) 상태로 자동 세팅"
        ])
    ]

    for idx, (title, items) in enumerate(faq_inquiry):
        left = Inches(0.8 + idx * 5.9)
        top = Inches(1.7)
        card = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.6), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.3)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=15, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n기능 사양:\n\n", size_pt=12, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=12, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 8: ADMIN MANUAL 1 - STREAMLIT AI DASHBOARD
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide8, BG_WARM)
    add_header(slide8, "05. 운영자 매뉴얼 - Streamlit AI 자동화 대시보드")

    admin_stream = [
        ("1단계: 알라딘 Open API 도서 연동", [
            "• `streamlit_admin.py` 관리자 대시보드 실행",
            "• 알라딘 TTB 키 연동으로 책 제목 검색",
            "• 책 표지 이미지 URL, 저자, 알라딘 URL 자동 추출"
        ]),
        ("2단계: Gemini AI 리뷰어 생성", [
            "• Google Gemini API `SINGLE_ROOM_SCHEMA` 호동",
            "• 350자+ 3단락 체계적 참여자 생생 후기 자동 생성",
            "• 리뷰어 직업/페르소나 및 별점, 좋아요 자동 매핑"
        ]),
        ("3단계: Supabase DB 원클릭 출판", [
            "• [Supabase 출판하기] 클릭 시 즉시 DB 전송",
            "• `rooms` 테이블 (모임 정보) 생성",
            "• `reviews` 테이블 (AI 리뷰 4개) 상호 외래키 매핑"
        ])
    ]

    for idx, (title, items) in enumerate(admin_stream):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n운영 절차:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 9: ADMIN MANUAL 2 - SUPABASE DB TABLES
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide9, BG_WARM)
    add_header(slide9, "05. 운영자 매뉴얼 - Supabase 데이터베이스 테이블 관리")

    db_tables = [
        ("rooms 테이블 (독서모임)", [
            "• status: recruiting(모집중), in_progress(진행중), completed(종료)",
            "• title, book_title, book_author, book_image_url, aladin_url",
            "• schedule_text, price_text, leader_name, leader_bio"
        ]),
        ("reviews 테이블 (생생 후기)", [
            "• room_id: rooms 테이블 외래키 매핑",
            "• author_name, author_role, rating, content, like_count",
            "• is_ai_generated: true (AI 생성 여부 플래그)"
        ]),
        ("inquiries 테이블 (1:1 문의)", [
            "• name, email, phone, inquiry_type, subject, content",
            "• status: pending(접수대기) -> replied(답변완료)",
            "• created_at: 문의 접수 타임스탬프"
        ])
    ]

    for idx, (title, items) in enumerate(db_tables):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n테이블 스키마 구조:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 10: SYSTEM INFRA & VERCEL DEPLOYMENT
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide10, BG_WARM)
    add_header(slide10, "06. 시스템 인프라 & Vercel 배포 가이드")

    infra_guide = [
        ("Vercel 프로덕션 빌드 설정", [
            "• Framework Preset: Vite",
            "• Build Command: npm run build",
            "• Output Directory: out (vercel.json에 세팅)",
            "• Rewrites: React Router SPA 라우팅 지원"
        ]),
        ("알라딘 이미지 CDN 차단 방지 Proxy", [
            "• Serverless Function: api/book-cover.js",
            "• 알라딘 커버 이미지 요청 시 referrerPolicy 해결",
            "• 엑박 방지 및 24시간 HTTP 캐싱 헤더 주입"
        ]),
        ("Supabase RLS & 보안 권한", [
            "• Read Access: anon(비로그인) 누구나 rooms/reviews 조회",
            "• Write Access: inquiries 테이블 INSERT 허용",
            "• service_role 키는 클라이언트에 절대 노출 안 됨"
        ])
    ]

    for idx, (title, items) in enumerate(infra_guide):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n인프라 아키텍처:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 11: VERCEL ENVIRONMENT VARIABLES (.ENV)
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide11, BG_WARM)
    add_header(slide11, "06. Vercel 환경 변수 (.env) 셋팅 가이드")

    # Table
    rows, cols = 4, 3
    table_shape = slide11.shapes.add_table(rows, cols, Inches(0.8), Inches(1.7), Inches(11.733), Inches(3.2))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(5.5)
    table.columns[2].width = Inches(3.033)

    headers = ["환경변수 Key (Vercel ENV)", "설정 값 (Value)", "설명 및 사용처"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_DARK
        p = cell.text_frame.paragraphs[0]
        add_run(p, h, size_pt=12, bold=True, color=TEXT_WHITE)

    env_data = [
        ("VITE_SUPABASE_URL", "https://fewzfqkqmdfbfqdtxmvf.supabase.co", "Supabase 백엔드 데이터베이스 API URL"),
        ("VITE_SUPABASE_ANON_KEY", "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9...", "Supabase Public Anon Key"),
        ("VITE_ALADIN_TTB_KEY", "ttbjwl17220625001", "알라딘 300권 Open API TTB Key")
    ]

    for row_idx, row_data in enumerate(env_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE
            p = cell.text_frame.paragraphs[0]
            add_run(p, cell_data, size_pt=11, bold=False, color=TEXT_BLACK)

    # Alert Card
    alert = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.7))
    alert.fill.solid()
    alert.fill.fore_color.rgb = RGBColor(253, 248, 245)
    alert.line.color.rgb = COLOR_BRICK

    tf_a = alert.text_frame
    tf_a.word_wrap = True
    tf_a.margin_left = tf_a.margin_top = Inches(0.25)

    p = tf_a.paragraphs[0]
    add_run(p, "⚠️ Vercel 배포 시 필독 주의사항 (Security & Redeploy Checklist)\n", size_pt=13, bold=True, color=COLOR_BRICK)

    p = tf_a.add_paragraph()
    add_run(p, "1. Vercel Dashboard [Project Settings] -> [Environment Variables]에 위 3개 변수를 등록하세요.\n2. VITE_ 접두사 변수가 추가/수정된 후에는 반드시 [Redeploy]를 진행해야 빌드 파일에 올바르게 적용됩니다.", size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 12: CLOSING & Q&A
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide12, BG_DARK)

    frame12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6), Inches(12.133), Inches(6.3))
    frame12.fill.background()
    frame12.line.color.rgb = COLOR_BRICK
    frame12.line.width = Pt(2.5)

    tb12 = slide12.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.933), Inches(4.5))
    tf12 = tb12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    add_run(p, "QUESTIONS & SUPPORT\n", size_pt=13, bold=True, color=COLOR_BRICK)

    p = tf12.add_paragraph()
    add_run(p, "Q & A 및 운영 문의 안내\n\n", size_pt=34, bold=True, color=TEXT_WHITE)

    p = tf12.add_paragraph()
    add_run(p, "• PM / 서비스 총괄: ", size_pt=14, bold=True, color=COLOR_BRICK)
    add_run(p, "이재우 (Jaewoo Lee)\n", size_pt=14, bold=False, color=TEXT_WHITE)

    p = tf12.add_paragraph()
    add_run(p, "• GitHub Repository: ", size_pt=14, bold=True, color=COLOR_BRICK)
    add_run(p, "https://github.com/JaewooLee-AI/questionity_main_v2.git\n", size_pt=14, bold=False, color=TEXT_WHITE)

    p = tf12.add_paragraph()
    add_run(p, "• 고객 문의 / 지원: ", size_pt=14, bold=True, color=COLOR_BRICK)
    add_run(p, "support@questionity.kr / jwlee@project.com\n\n", size_pt=14, bold=False, color=TEXT_WHITE)

    p = tf12.add_paragraph()
    add_run(p, "Thank you for using Questionity!", size_pt=18, bold=True, color=COLOR_BRICK)

    out_path = "/Users/jwlee/project/questionity_main_v2/Questionity_User_and_Admin_Manual.pptx"
    prs.save(out_path)
    print(f"SUCCESSFULLY GENERATED PPT: {out_path}")

if __name__ == "__main__":
    create_manual_presentation()
