import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # Colors
    BG_WARM = RGBColor(246, 245, 241)     # #F6F5F1
    BG_DARK = RGBColor(17, 17, 17)        # #111111
    BG_WHITE = RGBColor(255, 255, 255)
    TEXT_BLACK = RGBColor(20, 20, 20)     # Explicit Deep Black
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(90, 90, 90)
    COLOR_BRICK = RGBColor(140, 35, 24)   # #8C2318
    COLOR_OLIVE = RGBColor(59, 72, 54)    # #3B4836
    BORDER_GRAY = RGBColor(190, 185, 175)

    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_run(paragraph, text, size_pt=12, bold=False, color=TEXT_BLACK, font_name="Arial"):
        run = paragraph.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color
        return run

    def add_header(slide, title_text, category_text="QUESTIONITY HANDOVER"):
        # Header Box
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p0 = tf.paragraphs[0]
        add_run(p0, category_text.upper(), size_pt=11, bold=True, color=COLOR_BRICK)

        p1 = tf.add_paragraph()
        add_run(p1, title_text, size_pt=22, bold=True, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 1: COVER
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide1, BG_DARK)

    # Outer Frame
    frame = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6), Inches(12.133), Inches(6.3))
    frame.fill.background()
    frame.line.color.rgb = COLOR_BRICK
    frame.line.width = Pt(2.5)

    tb1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.933), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    add_run(p, "QUESTION + COMMUNITY PLATFORM\n", size_pt=13, bold=True, color=COLOR_BRICK)

    p = tf1.add_paragraph()
    add_run(p, "퀘스처니티 (Questionity)\n", size_pt=38, bold=True, color=TEXT_WHITE)

    p = tf1.add_paragraph()
    add_run(p, "1차 개발 완료 인수인계서 (Handover Document)\n\n", size_pt=22, bold=False, color=RGBColor(210, 210, 210))

    p = tf1.add_paragraph()
    add_run(p, "• 문서 목적: ", size_pt=13, bold=True, color=COLOR_BRICK)
    add_run(p, "퀘스처니티 시스템 IA, DB 환경변수, Vercel 배포 및 관리자/사용자 운영 매뉴얼\n", size_pt=13, bold=False, color=RGBColor(190, 190, 190))

    p = tf1.add_paragraph()
    add_run(p, "• 작성일자: ", size_pt=12, bold=True, color=COLOR_BRICK)
    add_run(p, "2026년 8월 14일   |   ", size_pt=12, bold=False, color=TEXT_WHITE)
    add_run(p, "• 작성자: ", size_pt=12, bold=True, color=COLOR_BRICK)
    add_run(p, "IT 프로젝트 매니저 (PM)", size_pt=12, bold=False, color=TEXT_WHITE)

    # =========================================================================
    # SLIDE 2: TABLE OF CONTENTS
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide2, BG_WARM)
    add_header(slide2, "목차 (Table of Contents)")

    toc_data = [
        ("01", "서비스 정보구조도 (IA)", "메인 서비스, 서브 페이지(/faq, /terms) 및 관리자 3계층 트리 구조"),
        ("02", "개발 환경 & 시스템 인프라", "React 18, Supabase DB, Python Streamlit Admin, Vercel 인프라"),
        ("03", "시스템 환경 변수 (.env) 가이드", "Supabase API Key, 알라딘 TTB Key, Vercel 환경 변수 및 배포 가이드"),
        ("04", "관리자 & 사용자 운영 매뉴얼", "Streamlit AI 리뷰 파이프라인, 알라딘 도서 연동, 1:1 문의 접수 처리")
    ]

    for idx, (num, title, desc) in enumerate(toc_data):
        col = idx % 2
        row = idx // 2
        left = Inches(0.8 + col * 5.9)
        top = Inches(1.7 + row * 2.6)

        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.6), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.3)
        tf_c.margin_top = Inches(0.3)

        p = tf_c.paragraphs[0]
        add_run(p, f"SECTION {num}\n", size_pt=10, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, title, size_pt=16, bold=True, color=TEXT_BLACK)

        p = tf_c.add_paragraph()
        add_run(p, f"\n{desc}", size_pt=12, bold=False, color=TEXT_MUTED)

    # =========================================================================
    # SLIDE 3: IA - MAIN SERVICE
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide3, BG_WARM)
    add_header(slide3, "01. 서비스 정보구조도 (IA) - 메인 서비스 (Home)")

    ia_main = [
        ("Depth 2: 히어로 (Hero)", [
            "• 브랜드 비전 (Question + Community)",
            "• 독서모임 둘러보기 / 결제방법 CTA"
        ]),
        ("Depth 2: 추천 독서모임 (Clubs)", [
            "• 모집중 / 진행중 / 종료 탭 필터",
            "• 📚 알라딘 바로가기 / 상세 모달"
        ]),
        ("Depth 2: 이렇게 시작해요", [
            "• 3단계 쿠폰 결제 흐름도",
            "• 2개 / 6개 / 10개 패키지 카드"
        ]),
        ("Depth 2: 오시는 길 (Location)", [
            "• 창경궁로 270 (혜화역 4번 출구)",
            "• 네이버 지도 길찾기 / 주소 복사"
        ])
    ]

    for idx, (title, items) in enumerate(ia_main):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.7)
        card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.8), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=13, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\nDepth 3 기능 목록:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 4: IA - SUB PAGES & ADMIN
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide4, BG_WARM)
    add_header(slide4, "01. 서비스 정보구조도 (IA) - 서브 페이지 & 관리자")

    ia_sub = [
        ("FAQ / 문의하기 (/faq)", [
            "• 실시간 질문 검색바",
            "• 카테고리 탭 (모임/결제/장소/환불)",
            "• Q&A 아코디언",
            "• 1:1 문의 폼 (inquiries DB 접수)"
        ]),
        ("이용약관 및 정책 (/terms, /privacy)", [
            "• 이용약관 (환불 규정 100% 명시)",
            "• 개인정보처리방침 (수집/보유/파기)",
            "• 개인정보 보호책임자 정보"
        ]),
        ("관리자 파이프라인 (Admin Streamlit)", [
            "• 도서 검색 & 알라딘 TTB 키 연동",
            "• Gemini AI 350자+ 리뷰 4개 자동 생성",
            "• Supabase rooms/reviews DB 원클릭 게시"
        ])
    ]

    for idx, (title, items) in enumerate(ia_sub):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n세부 구성 요소:\n\n", size_pt=12, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=12, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 5: TECH STACK & INFRA
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide5, BG_WARM)
    add_header(slide5, "02. 개발 환경 및 시스템 인프라 개요")

    stacks = [
        ("프론트엔드 (Frontend)", [
            "• Core: React 18, TypeScript, Vite",
            "• Styling: TailwindCSS (Ace Hotel 디자인)",
            "• Theme: #F6F5F1 웜캔버스, #111111 매트블랙",
            "• Layout: 100% 직사각형(rounded-none)"
        ]),
        ("백엔드 & DB (Backend & DB)", [
            "• Database: Supabase PostgreSQL",
            "• Main Tables: rooms, reviews, inquiries",
            "• Auth & Realtime: Supabase JS Client (v2)",
            "• Security: Row Level Security (RLS) 적용"
        ]),
        ("배포 & 인프라 (Infra & Pipeline)", [
            "• Web Hosting: Vercel Production",
            "• SPA Rewrites: vercel.json (outDir: out)",
            "• Serverless API: api/book-cover.js",
            "• Admin Pipeline: Streamlit + Gemini API"
        ])
    ]

    for idx, (title, items) in enumerate(stacks):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n주요 항목:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 6: ENV VARS
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide6, BG_WARM)
    add_header(slide6, "03. 시스템 환경 변수 (.env) 설정 가이드")

    # Table
    rows, cols = 4, 3
    table_shape = slide6.shapes.add_table(rows, cols, Inches(0.8), Inches(1.7), Inches(11.733), Inches(3.2))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(5.5)
    table.columns[2].width = Inches(3.033)

    headers = ["환경변수 Key", "설명 및 설정 값 예시", "비고 / 사용처"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_DARK
        p = cell.text_frame.paragraphs[0]
        add_run(p, h, size_pt=12, bold=True, color=TEXT_WHITE)

    env_data = [
        ("VITE_SUPABASE_URL", "https://fewzfqkqmdfbfqdtxmvf.supabase.co", "Supabase 프로젝트 접속 URL"),
        ("VITE_SUPABASE_ANON_KEY", "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9...", "Supabase Public Anon Key"),
        ("VITE_ALADIN_TTB_KEY", "ttbjwl17220625001", "알라딘 Open API TTB Key")
    ]

    for row_idx, row_data in enumerate(env_data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_WHITE
            p = cell.text_frame.paragraphs[0]
            add_run(p, cell_data, size_pt=11, bold=False, color=TEXT_BLACK)

    # Alert Card
    alert = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.7))
    alert.fill.solid()
    alert.fill.fore_color.rgb = RGBColor(253, 248, 245)
    alert.line.color.rgb = COLOR_BRICK

    tf_a = alert.text_frame
    tf_a.word_wrap = True
    tf_a.margin_left = tf_a.margin_top = Inches(0.25)

    p = tf_a.paragraphs[0]
    add_run(p, "⚠️ Vercel 배포 및 보안 주의사항 (Security Checklist)\n", size_pt=13, bold=True, color=COLOR_BRICK)

    p = tf_a.add_paragraph()
    add_run(p, "1. Vercel Dashboard [Project Settings] -> [Environment Variables]에 위 3가지 변수를 반드시 등록해야 합니다.\n2. VITE_ 접두사가 포함된 변수는 클라이언트에 공개되므로, Supabase service_role 키는 절대 등록하지 마십시오.", size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 7: ADMIN MANUAL 1
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide7, BG_WARM)
    add_header(slide7, "04. 관리자(Admin) 매뉴얼 - 독서방 생성 & AI 파이프라인")

    admin1 = [
        ("1단계: 독서방 기획 정보 입력", [
            "• Streamlit 대시보드 접속",
            "• 모임 제목, 책 제목, 저자, 모임 일시 입력",
            "• 알라딘 TTB 키로 도서 URL 자동 수집"
        ]),
        ("2단계: Gemini AI 리뷰 자동 생성", [
            "• Gemini API SINGLE_ROOM_SCHEMA 작동",
            "• 350자+ 3단락 체계적 참여자 생생 후기 작성",
            "• 리뷰어 페르소나 및 좋아요 수 자동 설정"
        ]),
        ("3단계: Supabase DB 원클릭 게시", [
            "• [Supabase 출판하기] 버튼 클릭",
            "• rooms 테이블 및 reviews 테이블 상호 매핑",
            "• 메인 웹 서비스 추천 독서모임 즉시 노출"
        ])
    ]

    for idx, (title, items) in enumerate(admin1):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n세부 절차:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 8: ADMIN MANUAL 2
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide8, BG_WARM)
    add_header(slide8, "04. 관리자(Admin) 매뉴얼 - 고객 문의 & 데이터 관리")

    admin2 = [
        ("고객 1:1 문의 관리 (inquiries 테이블)", [
            "• 접수 경로: 웹사이트 /faq#contact 폼 제출",
            "• 수집 데이터: 성함, 이메일, 연락처, 문의유형, 제목, 내용",
            "• 상태 관리: pending(접수대기) -> replied(답변완료)"
        ]),
        ("알라딘 도서 이미지 CDN 차단 관리", [
            "• Vercel Serverless Function: api/book-cover.js",
            "• referrerPolicy='no-referrer' 설정으로 이미지 엑박 방지",
            "• 실시간 캐싱(Cache-Control 24h) 적용"
        ])
    ]

    for idx, (title, items) in enumerate(admin2):
        left = Inches(0.8 + idx * 5.9)
        top = Inches(1.7)
        card = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.6), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.3)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=15, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n운용 규칙:\n\n", size_pt=12, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=12, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 9: USER MANUAL 1
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide9, BG_WARM)
    add_header(slide9, "04. 사용자(User) 매뉴얼 - 독서방 탐색 & 알라딘 연동")

    user1 = [
        ("1. 독서모임 탐색 및 필터링", [
            "• 메인 화면 #clubs 이동",
            "• 모집중 / 진행중 / 종료 탭 필터 제공",
            "• 모임 카드 직사각형 서식 적용"
        ]),
        ("2. 알라딘 도서 바로가기 (📚 알라딘 ↗)", [
            "• 카드 및 상세 모달에서 바로가기 버튼",
            "• 정제된 책 제목 기반 구매 페이지 이동",
            "• 도서 정보 상세 확인 가능"
        ]),
        ("3. 독서모임 상세 모달 팝업", [
            "• 카드 클릭 시 모임 소개 & 도서 팝업",
            "• 클럽장 프로필 및 이력 확인",
            "• AI 생성 참여자 실전 후기 조회"
        ])
    ]

    for idx, (title, items) in enumerate(user1):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n사용자 흐름:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 10: USER MANUAL 2
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide10, BG_WARM)
    add_header(slide10, "04. 사용자(User) 매뉴얼 - 쿠폰 결제, 오프라인 아지트 & FAQ")

    user2 = [
        ("쿠폰 결제 안내 (#how-it-works)", [
            "• 2개(5%), 6개(15% 추천), 10개(20% 할인)",
            "• 계좌이체, 카드등록, NPay, KakaoPay 지원",
            "• 독서방 결제 시 [쿠폰 1장 차감] 신청"
        ]),
        ("오프라인 아지트 안내 (#location)", [
            "• 주소: 서울특별시 종로구 창경궁로 270",
            "• 위치: 4호선 혜화역 4번 출구 (도보 3분)",
            "• 기능: [도로명 주소 복사] & [네이버 지도 길찾기]"
        ]),
        ("FAQ & 1:1 문의 접수 (/faq)", [
            "• 실시간 질문/답변 키워드 검색바",
            "• 4개 카테고리 탭 및 아코디언 Q&A",
            "• 1:1 문의 제출 시 Supabase DB 자동 접수"
        ])
    ]

    for idx, (title, items) in enumerate(user2):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.7)
        card = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.75), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = BORDER_GRAY

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = Inches(0.25)

        p = tf_c.paragraphs[0]
        add_run(p, title, size_pt=14, bold=True, color=COLOR_BRICK)

        p = tf_c.add_paragraph()
        add_run(p, "\n서비스 기능:\n\n", size_pt=11, bold=True, color=TEXT_BLACK)
        for item in items:
            p_item = tf_c.add_paragraph()
            add_run(p_item, item, size_pt=11, bold=False, color=TEXT_BLACK)

    # =========================================================================
    # SLIDE 11: CLOSING & Q&A
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide11, BG_DARK)

    frame11 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6), Inches(12.133), Inches(6.3))
    frame11.fill.background()
    frame11.line.color.rgb = COLOR_BRICK
    frame11.line.width = Pt(2.5)

    tb11 = slide11.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.933), Inches(4.5))
    tf11 = tb11.text_frame
    tf11.word_wrap = True

    p = tf11.paragraphs[0]
    add_run(p, "QUESTIONS & ANSWERS\n", size_pt=13, bold=True, color=COLOR_BRICK)

    p = tf11.add_paragraph()
    add_run(p, "Q & A 및 인수인계 담당자 연락처\n\n", size_pt=34, bold=True, color=TEXT_WHITE)

    p = tf11.add_paragraph()
    add_run(p, "• PM / Tech Lead: ", size_pt=14, bold=True, color=COLOR_BRICK)
    add_run(p, "이재우 (Jaewoo Lee)\n", size_pt=14, bold=False, color=TEXT_WHITE)

    p = tf11.add_paragraph()
    add_run(p, "• GitHub Repository: ", size_pt=14, bold=True, color=COLOR_BRICK)
    add_run(p, "https://github.com/JaewooLee-AI/questionity_main.git\n", size_pt=14, bold=False, color=TEXT_WHITE)

    p = tf11.add_paragraph()
    add_run(p, "• Support Email: ", size_pt=14, bold=True, color=COLOR_BRICK)
    add_run(p, "support@questionity.kr / jwlee@project.com\n\n", size_pt=14, bold=False, color=TEXT_WHITE)

    p = tf11.add_paragraph()
    add_run(p, "Thank you for your partnership!", size_pt=18, bold=True, color=COLOR_BRICK)

    out_path = "/Users/jwlee/project/questionity_main/Questionity_Project_Handover.pptx"
    prs.save(out_path)
    print(f"RE-GENERATED SUCCESS: {out_path}")

if __name__ == "__main__":
    create_presentation()
